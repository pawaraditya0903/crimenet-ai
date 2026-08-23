import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

C_BG = RGBColor(10, 16, 31)
C_CARD = RGBColor(15, 23, 42)
C_BLUE = RGBColor(56, 189, 248)
C_GREEN = RGBColor(52, 211, 153)
C_RED = RGBColor(239, 68, 68)
C_WHITE = RGBColor(255, 255, 255)
C_GRAY = RGBColor(148, 163, 184)

def create_base_slide(title, subtitle="CRIMENET AI — ADITYA PAWAR PLATFORM"):
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_sub = tf.paragraphs[0]
    p_sub.text = subtitle.upper()
    p_sub.font.size = Pt(10)
    p_sub.font.bold = True
    p_sub.font.color.rgb = C_BLUE

    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE

    return slide

# SLIDE 1: Title
s1 = create_base_slide("CRIMENET AI — AUTONOMOUS CRIMINAL GRAPH & FORENSICS", "EXECUTIVE ENGINEERING PORTFOLIO")
tb = s1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Lead Architect: Aditya Pawar | Full-Stack & AI/ML Systems Engineer"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = C_BLUE

p2 = tf.add_paragraph()
p2.text = "\n• Sub-Second Multi-Hop Graph Traversal across 50,000+ logs (14 days -> 420ms)."
p2.text += "\n• Mathematical Kingpin Discovery via PageRank (0.0847) & Betweenness Centrality (0.312)."
p2.text += "\n• Hardware IMEI Co-Location Engine solving disposable burner SIM swaps."
p2.text += "\n• Multi-Hop Blockchain Crypto Hawala Tracer ($2.45M USDT) with Section 65B SHA-256 Hashes."
p2.text += "\n• Graph-RAG Zero-Hallucination AI Copilot using Google Gemini 1.5 Pro & Web Speech Synthesis."
p2.font.size = Pt(14)
p2.font.color.rgb = C_WHITE

# SLIDE 2: Architecture
s2 = create_base_slide("FULL-STACK SYSTEM DESIGN & ARCHITECTURE MATRIX", "TECHNICAL SPECIFICATIONS")
cards_data = [
    ("FRONTEND UI", "React 18, TypeScript Strict Mode, Cytoscape.js 60fps Physics, 2D Radar Canvas, Web Speech Synthesis.", C_BLUE),
    ("BACKEND ASGI", "Python 3.11 FastAPI, Asynchronous Non-Blocking Workers, Socket.IO WebSockets, ReportLab PDF Engine.", C_GREEN),
    ("GRAPH ML", "NetworkX PageRank & Louvain Modularity, Scikit-Learn Isolation Forest, IMEI Co-Location Fingerprinting.", C_BLUE),
    ("GRAPH-RAG AI", "Google Gemini 1.5 Pro, 2-Hop Grounded Subgraph Context, Section 65B Evidence Certification.", C_GREEN)
]
for idx, (head, desc, color) in enumerate(cards_data):
    x = Inches(0.8 + idx * 2.95)
    card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.8), Inches(4.2))
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD
    card.line.color.rgb = color
    ctf = card.text_frame
    ctf.word_wrap = True
    p1 = ctf.paragraphs[0]
    p1.text = head
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = color
    p2 = ctf.add_paragraph()
    p2.text = f"\n{desc}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = C_WHITE

# SLIDE 3: Candidate Profile
s3 = create_base_slide("WHY HIRE ADITYA PAWAR: FULL-STACK & AI/ML ENGINEER", "CANDIDATE CREDENTIALS")
tb3 = s3.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
tf3 = tb3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = "Key Engineering Strengths & Production Readiness:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = C_BLUE

skills = [
    "✓ Full-Stack Engineering: Expert in Python (FastAPI/ASGI), React 18, TypeScript, and 60fps Canvas physics.",
    "✓ Applied Machine Learning: Hands-on mathematical mastery of Graph Theory (PageRank, Louvain) and Anomaly Detection.",
    "✓ Production Scalability: Deep understanding of distributed stream processing (Kafka/Flink), Neo4j graphs, and LOD rendering.",
    "✓ High ROI Impact: Proven track record of reducing complex analysis workflows by 90%."
]
for sk in skills:
    p_sk = tf3.add_paragraph()
    p_sk.text = f"\n{sk}"
    p_sk.font.size = Pt(13)
    p_sk.font.color.rgb = C_WHITE

prs.save("c:\\Users\\Aditya\\Downloads\\SIH 2026\\CrimeNet_AI_Aditya_Pawar_Master_Deck.pptx")
print("PowerPoint presentation generated: CrimeNet_AI_Aditya_Pawar_Master_Deck.pptx")
