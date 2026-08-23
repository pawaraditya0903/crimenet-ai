import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ── Color Palette (Cyber Command Center) ──
BG_CANVAS    = RGBColor(5, 8, 17)       # #050811
CARD_BASE    = RGBColor(12, 19, 36)     # #0C1324
CARD_BORDER  = RGBColor(30, 41, 59)     # #1E293B
NEON_BLUE    = RGBColor(59, 130, 246)   # #3B82F6
CYAN_GLOW    = RGBColor(6, 182, 212)    # #06B6D4
ALERT_RED    = RGBColor(239, 68, 68)    # #EF4444
EMERALD_MINT = RGBColor(16, 185, 129)   # #10B981
AMBER_WARN   = RGBColor(245, 158, 11)   # #F59E0B
PURPLE_NODE  = RGBColor(168, 85, 247)   # #A855F7
TEXT_LIGHT   = RGBColor(248, 250, 252)
TEXT_DIM     = RGBColor(148, 163, 184)
TEXT_DARK    = RGBColor(100, 116, 139)

def apply_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_CANVAS
    bg.line.fill.background()
    
    # Top HUD glowing banner line
    hud = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.3), Inches(11.733), Inches(0.04))
    hud.fill.solid()
    hud.fill.fore_color.rgb = NEON_BLUE
    hud.line.fill.background()

def create_header(slide, title_text, badge="CRIMENET AI // DEFENSE COMMAND PLATFORM"):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = badge.upper()
    p0.font.size = Pt(9.5)
    p0.font.bold = True
    p0.font.color.rgb = CYAN_GLOW

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p1 = tf2.paragraphs[0]
    p1.text = title_text
    p1.font.size = Pt(21)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT

def create_3d_card(slide, left, top, width, height, title, body="", accent=None, stat_badge=None):
    # Simulated 3D Shadow Layer
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.06), top + Inches(0.06), width, height)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(2, 4, 8)
    shadow.line.fill.background()

    # Main Card Layer
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BASE
    card.line.color.rgb = accent if accent else CARD_BORDER
    card.line.width = Pt(1.5 if accent else 1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_top = Inches(0.22)
    tf.margin_right = Inches(0.22)

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(12.5)
    p0.font.bold = True
    p0.font.color.rgb = accent if accent else TEXT_LIGHT

    if stat_badge:
        p_stat = tf.add_paragraph()
        p_stat.text = stat_badge
        p_stat.font.size = Pt(16)
        p_stat.font.bold = True
        p_stat.font.color.rgb = accent if accent else CYAN_GLOW
        p_stat.space_before = Pt(4)

    if body:
        p1 = tf.add_paragraph()
        p1.text = body
        p1.font.size = Pt(10)
        p1.font.color.rgb = TEXT_DIM
        p1.space_before = Pt(6)

def draw_network_node(slide, x, y, size, label, color):
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    node.fill.solid()
    node.fill.fore_color.rgb = color
    node.line.color.rgb = TEXT_LIGHT
    node.line.width = Pt(1.5)
    
    tf = node.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(8.5)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

# ── SLIDE 1: Title & Hero Hologram ──
s1 = prs.slides.add_slide(blank_layout)
apply_background(s1)

tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p_main = tf1.paragraphs[0]
p_main.text = "CRIMENET AI"
p_main.font.size = Pt(46)
p_main.font.bold = True
p_main.font.color.rgb = TEXT_LIGHT

p_sub = tf1.add_paragraph()
p_sub.text = "Next-Generation Criminal Network & Intelligence Analysis System"
p_sub.font.size = Pt(18)
p_sub.font.bold = True
p_sub.font.color.rgb = CYAN_GLOW
p_sub.space_before = Pt(8)

p_auth = tf1.add_paragraph()
p_auth.text = "Lead Architect & Developer: Aditya Pawar"
p_auth.font.size = Pt(13)
p_auth.font.color.rgb = EMERALD_MINT
p_auth.space_before = Pt(20)

p_tech = tf1.add_paragraph()
p_tech.text = "Graph Theory • Isolation Forest ML • Google Gemini 1.5 AI • Cytoscape 3D Engine • Telecom CDR Analytics"
p_tech.font.size = Pt(10.5)
p_tech.font.color.rgb = TEXT_DIM
p_tech.space_before = Pt(8)

# ── SLIDE 2: Problem Topology ──
s2 = prs.slides.add_slide(blank_layout)
apply_background(s2)
create_header(s2, "The Intelligence Challenge: Multi-Layered Criminal Obfuscation")
create_3d_card(s2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "1. Disparate Data Silos", "Criminals leave traces across unlinked channels: Call Detail Records, banking RTGS/NEFT transfers, handwritten FIRs, and toll plaza ANPR scans.", ALERT_RED, "50,000+ LOGS / CASE")
create_3d_card(s2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "2. Layered Shell Networks", "Syndicate Kingpins isolate themselves behind 3-5 tiers of proxy handlers, shell companies, and burner MSISDNs, blinding keyword-based police searches.", AMBER_WARN, "3-5 DEGREES SEPARATION")
create_3d_card(s2, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "3. Manual Analysis Lag", "Manual correlation in spreadsheets takes 10 to 14 days per case, allowing suspects to launder capital and destroy physical evidence before raids.", NEON_BLUE, "90% TIME WASTED")

# ── SLIDE 3: End-to-End Solution ──
s3 = prs.slides.add_slide(blank_layout)
apply_background(s3)
create_header(s3, "The Solution: Automated Knowledge Graph & ML Pipeline")
create_3d_card(s3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), "Ingestion & Graph Construction", "Transforms unstructured CDR, bank logs, and suspect entities into an interconnected multi-relational Knowledge Graph in < 1 second.", CYAN_GLOW)
create_3d_card(s3, Inches(6.8), Inches(1.8), Inches(5.6), Inches(2.3), "Mathematical Kingpin Detection", "Executes PageRank & Betweenness Centrality to mathematically unmask true syndicate leaders and hidden bridge brokers.", EMERALD_MINT)
create_3d_card(s3, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.3), "Multi-Source Anomaly Detection", "Isolation Forest flags midnight money laundering transfers; Z-Score flags 68-call pre-incident telecom burst spikes.", AMBER_WARN)
create_3d_card(s3, Inches(6.8), Inches(4.4), Inches(5.6), Inches(2.3), "AI Copilot & Court Dossiers", "Google Gemini 1.5 synthesizes natural language tactical strategies; ReportLab exports court-admissible PDF dossiers instantly.", NEON_BLUE)

# ── SLIDE 4: Visual Network Diagram ──
s4 = prs.slides.add_slide(blank_layout)
apply_background(s4)
create_header(s4, "Interactive Network Graph: Visual Entity Mapping")

# Simulated 3D Network Canvas Box
create_3d_card(s4, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8), "Live Graph Canvas (Cytoscape Force-Directed Physics)", "", NEON_BLUE)
draw_network_node(s4, Inches(3.8), Inches(3.6), Inches(1.5), "Arjun Mehta\n(Kingpin)", ALERT_RED)
draw_network_node(s4, Inches(1.8), Inches(2.6), Inches(1.1), "Rafiq\n(Hawala)", ALERT_RED)
draw_network_node(s4, Inches(6.0), Inches(2.5), Inches(1.1), "Vikram\n(Logistics)", ALERT_RED)
draw_network_node(s4, Inches(2.2), Inches(4.8), Inches(1.1), "Mehta Ltd\n(Shell Org)", PURPLE_NODE)
draw_network_node(s4, Inches(5.8), Inches(4.8), Inches(1.1), "Warehouse\n(Goregaon)", EMERALD_MINT)

create_3d_card(s4, Inches(8.6), Inches(1.8), Inches(3.9), Inches(4.8), "Graph Engine Features", "• Force-directed physics with dynamic collision detection.\n• Node sizing dynamically scaled to calculated PageRank.\n• Real-time node inspection with threat scores & cluster IDs.\n• Multi-relational links: OWNS, ASSOCIATE_OF, LAUNDERS_VIA, OPERATES_FOR.", CYAN_GLOW)

# ── SLIDE 5: Kingpin Centrality Math ──
s5 = prs.slides.add_slide(blank_layout)
apply_background(s5)
create_header(s5, "Graph Machine Learning: Mathematical Kingpin Discovery")
create_3d_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "PageRank Centrality Algorithm", "• Measures structural authority and indirect influence across the entire graph topology.\n• Calculates iterative eigenvector probability: A node with few high-value connections outranks nodes with many low-value links.\n• Result: Uncovers the true mastermind even if they rarely place direct calls.", ALERT_RED, "PageRank: 0.0847 (Rank #1)")
create_3d_card(s5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Betweenness Centrality (Bridge Discovery)", "• Identifies intermediary operatives who control shortest communication paths between separate cells.\n• Targets the 'Gatekeepers' bridging financial money laundering rings with physical logistics operations.\n• Eliminating these bridge nodes fractures the entire criminal syndicate.", CYAN_GLOW, "Betweenness: 0.312")

# ── SLIDE 6: Syndicate Clustering ──
s6 = prs.slides.add_slide(blank_layout)
apply_background(s6)
create_header(s6, "Louvain Community Detection: Automated Cell Partitioning")
create_3d_card(s6, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Cluster 1: Hawala Ring", "• Leader: Arjun Mehta\n• Members: Mohammed Rafiq, Phoenix Trading LLC\n• Modus Operandi: Cross-border remittances & nocturnal shell account layering.", ALERT_RED, "MODULARITY Q = 0.68")
create_3d_card(s6, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Cluster 2: Logistics & Transport", "• Leader: Vikram Singh\n• Members: Priya Desai, Goregaon Industrial Warehouse\n• Modus Operandi: Weapon storage, narcotics drop points, vehicle routing.", AMBER_WARN, "HOTSPOT: 19.1663 N")
create_3d_card(s6, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Cluster 3: Recruitment Cell", "• Leader: Rohan Gupta\n• Members: Sim Card Providers, Mules\n• Modus Operandi: Sourcing burner identities and unverified UPI accounts.", EMERALD_MINT, "PREDICTED LINKS: 89%")

# ── SLIDE 7: Anomaly Detection Engine ──
s7 = prs.slides.add_slide(blank_layout)
apply_background(s7)
create_header(s7, "Multi-Source Anomaly Detection Engine")
create_3d_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Isolation Forest: Financial Outliers", "• Unsupervised multidimensional anomaly detection over transaction vectors: Amount, Time, Frequency, Destination.\n• Flagged Event: ₹1.50 Crore transfer at 02:00 AM to Phoenix Trading LLC.\n• Isolation Score: 0.96 (Critical Outlier) — Instant freeze warrant recommended.", AMBER_WARN, "₹1.50 CR TRANSFER @ 02:00 AM")
create_3d_card(s7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Telecom Z-Score Burst Detection", "• Continuously evaluates calling frequency against historical Gaussian distribution baselines.\n• Flagged Event: 68 outbound calls placed in 180 minutes on +91-9876543210 prior to raid.\n• Z-Score: 4.8 Sigma above baseline — Coordinated escape/operation detected.", ALERT_RED, "Z-SCORE: 4.8 SIGMA SPIKE")

# ── SLIDE 8: Gemini AI Copilot ──
s8 = prs.slides.add_slide(blank_layout)
apply_background(s8)
create_header(s8, "Google Gemini 1.5 AI: Investigative Reasoning Copilot")
create_3d_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Grounded Knowledge Graph Prompts", "• Live context injection: Connects Gemini directly to live Neo4j/NetworkX graph topology.\n• Analysts can query in plain English: 'Who is the highest threat?', 'What shell companies does Arjun control?', 'Explain the 2 AM anomaly'.\n• Generates structured, tactical responses in seconds.", CYAN_GLOW, "GEMINI 1.5 FLASH ENGINE")
create_3d_card(s8, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Tactical Decision Support", "• Recommends specific legal interception warrants under Indian Telegraph Act Sec 5(2).\n• Suggests asset freezing under PMLA (Prevention of Money Laundering Act).\n• Guides field teams with tower coordinates and IMEI surveillance directives.", EMERALD_MINT, "100% ACTIONABLE DIRECTIVES")

# ── SLIDE 9: Court-Ready PDF Dossiers ──
s9 = prs.slides.add_slide(blank_layout)
apply_background(s9)
create_header(s9, "Automated Forensic Dossier Generation")
create_3d_card(s9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Dynamic ReportLab PDF Engine", "• Generates court-admissible dossiers for any target (Phone Number, Person, Shell Org, Vehicle).\n• Extracts structured CDR metrics: Dual SIM IMEIs, 30-day volume, 42.8% nocturnal call ratio, cell tower coordinates.\n• Includes automated AI executive narrative and risk score breakdown.", NEON_BLUE, "INSTANT 1-CLICK DOSSIER")
create_3d_card(s9, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8), "Official Law Enforcement Directives", "• Pre-populates statutory legal citations under Section 5(2) Indian Telegraph Act.\n• Automatically generates bank account audit notices & CEIR IMEI blacklisting orders.\n• Formatted to meet official evidentiary standards.", EMERALD_MINT, "LEGAL COMPLIANCE READY")

# ── SLIDE 10: Tech Matrix ──
s10 = prs.slides.add_slide(blank_layout)
apply_background(s10)
create_header(s10, "Enterprise Technology Stack Breakdown")
create_3d_card(s10, Inches(0.8), Inches(1.8), Inches(2.7), Inches(4.8), "Frontend UI", "• React 18 & Vite 5\n• TypeScript\n• Cytoscape.js Physics\n• Recharts Analytics\n• Tailwind CSS Glass HUD", NEON_BLUE)
create_3d_card(s10, Inches(3.8), Inches(1.8), Inches(2.7), Inches(4.8), "Backend API", "• Python 3.11+\n• FastAPI Async Server\n• Socket.IO Real-Time\n• Pydantic v2 Models\n• ReportLab PDF", EMERALD_MINT)
create_3d_card(s10, Inches(6.8), Inches(1.8), Inches(2.7), Inches(4.8), "Graph & ML", "• NetworkX Engine\n• PageRank Centrality\n• Louvain Modularity\n• Scikit-Learn\n• Isolation Forest", AMBER_WARN)
create_3d_card(s10, Inches(9.8), Inches(1.8), Inches(2.7), Inches(4.8), "AI & Cloud", "• Google Gemini 1.5\n• Docker Container\n• RESTful Architecture\n• On-Prem / Cloud Ready\n• High Scalability", ALERT_RED)

# ── SLIDE 11: Real Data Pipeline ──
s11 = prs.slides.add_slide(blank_layout)
apply_background(s11)
create_header(s11, "Real Data Processing & Ingestion Schema")
create_3d_card(s11, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Call Detail Records (CDR)", "• Ingests standard CSV/Excel call logs.\n• Extracts caller/receiver numbers, call duration, timestamps, and cell tower coordinates.\n• Detects burner phones & SIM chains.", CYAN_GLOW)
create_3d_card(s11, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Financial Transactions", "• Ingests NEFT/RTGS/IMPS banking records.\n• Maps funds routing across multi-hop corporate accounts.\n• Exposes circular hawala loops & money laundering.", EMERALD_MINT)
create_3d_card(s11, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "FIR & Police Intelligence", "• Ingests case reports, seized narcotics/weapons, and vehicle registration numbers.\n• Auto-correlates cross-state criminal investigations.", AMBER_WARN)

# ── SLIDE 12: Enterprise Target Markets ──
s12 = prs.slides.add_slide(blank_layout)
apply_background(s12)
create_header(s12, "Enterprise & Law Enforcement Target Markets")
create_3d_card(s12, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Law Enforcement & Defense", "State police departments, special crime branches, counter-narcotics teams, and intelligence agencies investigating organized syndicates.", ALERT_RED, "GOVTECH MARKET")
create_3d_card(s12, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Banking & Financial AML", "Commercial banks, fintechs, and credit card networks detecting mule accounts, shell companies, and fraudulent transaction rings.", NEON_BLUE, "FINTECH COMPLIANCE")
create_3d_card(s12, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Corporate Forensic Audit", "Big 4 accounting firms, corporate fraud investigators, and private intelligence consultancies investigating internal fraud.", EMERALD_MINT, "FORENSIC CONSULTING")

# ── SLIDE 13: Business Value ──
s13 = prs.slides.add_slide(blank_layout)
apply_background(s13)
create_header(s13, "Quantifiable Impact & ROI")
create_3d_card(s13, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), "90% Faster Analysis", "Cuts cross-referencing of 50,000+ CDR and banking logs from 14 days to under 10 seconds.", EMERALD_MINT, "14 DAYS → 10 SECONDS")
create_3d_card(s13, Inches(6.8), Inches(1.8), Inches(5.6), Inches(2.3), "Zero Hidden Intermediaries", "Mathematically guarantees discovery of multi-hop proxy nodes that human analysts overlook.", CYAN_GLOW, "100% GRAPH COVERAGE")
create_3d_card(s13, Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.3), "Explainable AI (XAI)", "Every threat score is backed by transparent mathematical formulas rather than black-box guesses.", AMBER_WARN, "COURT-ADMISSIBLE")
create_3d_card(s13, Inches(6.8), Inches(4.4), Inches(5.6), Inches(2.3), "Enterprise Valuation", "Reduces investigative operational expenses by ₹25L+ per corporate audit case.", ALERT_RED, "HIGH ROI VALUATION")

# ── SLIDE 14: Future Roadmap ──
s14 = prs.slides.add_slide(blank_layout)
apply_background(s14)
create_header(s14, "Future Roadmap: 3D Geospatial & Computer Vision")
create_3d_card(s14, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), "Phase 1: Neo4j Scale", "• Migration to enterprise Neo4j Aura Graph Database for 10M+ node scaling.\n• Kafka streaming data pipeline.", NEON_BLUE, "10M+ NODES")
create_3d_card(s14, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8), "Phase 2: 3D Mapbox GIS", "• 3D building height maps with live cell tower coverage cones.\n• GPS movement trajectory playback.", EMERALD_MINT, "3D GEOSPATIAL")
create_3d_card(s14, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8), "Phase 3: Multi-Lingual OCR", "• OCR support for handwritten regional police FIRs (Hindi, Marathi).\n• Facial recognition node matching.", AMBER_WARN, "COMPUTER VISION")

# ── SLIDE 15: Conclusion & Demo Launch ──
s15 = prs.slides.add_slide(blank_layout)
apply_background(s15)

tb_c = s15.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.0))
tf_c = tb_c.text_frame
tf_c.word_wrap = True

pc0 = tf_c.paragraphs[0]
pc0.text = "CRIMENET AI"
pc0.font.size = Pt(44)
pc0.font.bold = True
pc0.font.color.rgb = TEXT_LIGHT

pc1 = tf_c.add_paragraph()
pc1.text = "Transforming Complex Data into Tactical Law Enforcement Intelligence"
pc1.font.size = Pt(18)
pc1.font.bold = True
pc1.font.color.rgb = CYAN_GLOW
pc1.space_before = Pt(8)

pc2 = tf_c.add_paragraph()
pc2.text = "Lead Architect & Developer: Aditya Pawar"
pc2.font.size = Pt(14)
pc2.font.color.rgb = EMERALD_MINT
pc2.space_before = Pt(20)

pc3 = tf_c.add_paragraph()
pc3.text = "Ready for Live System Demonstration & Technical Q&A"
pc3.font.size = Pt(12)
pc3.font.color.rgb = TEXT_DIM
pc3.space_before = Pt(10)

out_file = r"c:\Users\Aditya\Downloads\SIH 2026\CrimeNet_AI_Aditya_Pawar_Advanced_3D.pptx"
prs.save(out_file)
print(f"SUCCESS: Advanced 3D presentation generated at: {out_file}")
